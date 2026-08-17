#!/usr/bin/env python3
"""Audit an image-only PPTX delivery package.

This is intentionally separate from audit_pptx_editability.py. A valid
image-only deck should look flattened in PPTX structure: one full-slide picture
per slide and no required editable text boxes.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_TOLERANCE = 25000
ASPECT_RATIO_TOLERANCE = 0.03


def _load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _resolve(base: Path, value: str | None) -> Path | None:
    if not value:
        return None
    path = Path(value)
    return path if path.is_absolute() else base / path


def _aspect_issue(label: str, path: Path, target_ratio: float) -> tuple[dict, str | None]:
    item = {"label": label, "path": str(path), "exists": path.exists()}
    if not path.exists():
        return item, f"{label}: image does not exist: {path}"
    try:
        with Image.open(path) as image:
            width, height = image.size
    except Exception as exc:  # pragma: no cover - defensive reporting
        return item, f"{label}: image is not readable: {exc}"

    ratio = width / max(height, 1)
    item.update({"width": width, "height": height, "aspect_ratio": round(ratio, 4)})
    if abs(ratio - target_ratio) / target_ratio > ASPECT_RATIO_TOLERANCE:
        return (
            item,
            f"{label}: image aspect ratio {width}:{height} ({ratio:.3f}) "
            f"does not match PPTX ratio {target_ratio:.3f}; refusing stretched full-slide packaging",
        )
    return item, None


def audit_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    slides = []
    total_pictures = 0
    total_textboxes = 0
    issues: list[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        pictures = []
        textboxes = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures.append(
                    {
                        "left": int(shape.left),
                        "top": int(shape.top),
                        "width": int(shape.width),
                        "height": int(shape.height),
                    }
                )
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                textboxes += 1

        total_pictures += len(pictures)
        total_textboxes += textboxes

        full_slide_pictures = [
            pic
            for pic in pictures
            if abs(pic["left"]) <= EMU_TOLERANCE
            and abs(pic["top"]) <= EMU_TOLERANCE
            and abs(pic["width"] - width) <= EMU_TOLERANCE
            and abs(pic["height"] - height) <= EMU_TOLERANCE
        ]
        if len(pictures) != 1:
            issues.append(f"slide {slide_index}: expected exactly 1 picture, found {len(pictures)}")
        if len(full_slide_pictures) != 1:
            issues.append(f"slide {slide_index}: picture does not cover the full slide")
        if textboxes:
            issues.append(f"slide {slide_index}: expected no editable text boxes, found {textboxes}")

        slides.append(
            {
                "slide_index": slide_index,
                "pictures": len(pictures),
                "full_slide_pictures": len(full_slide_pictures),
                "textboxes": textboxes,
                "picture_boxes": pictures,
            }
        )

    return {
        "pptx": str(path),
        "slide_count": len(prs.slides),
        "slide_width_emu": width,
        "slide_height_emu": height,
        "total_pictures": total_pictures,
        "total_textboxes": total_textboxes,
        "slides": slides,
        "issues": issues,
    }


def audit_package(
    pptx: Path,
    spec: Path | None,
    render_manifest: Path | None,
    deck_json: Path | None,
    min_slides: int,
) -> dict:
    report = audit_pptx(pptx)
    issues = list(report["issues"])
    target_ratio = report["slide_width_emu"] / max(report["slide_height_emu"], 1)
    source_images: list[dict] = []
    checked_paths: set[Path] = set()

    if report["slide_count"] < min_slides:
        issues.append(f"slide_count below {min_slides}")

    if spec:
        spec_data = _load_json(spec)
        spec_slides = spec_data.get("slides", [])
        if len(spec_slides) != report["slide_count"]:
            issues.append(
                f"slide_spec count {len(spec_slides)} does not match pptx slide_count {report['slide_count']}"
            )
        for slide in spec_slides:
            image_path = slide.get("rendered_image")
            if not image_path:
                issues.append(f"slide {slide.get('slide_id', '?')}: missing rendered_image")
                continue
            resolved = _resolve(spec.parent, image_path)
            if resolved and resolved not in checked_paths:
                checked_paths.add(resolved)
                item, issue = _aspect_issue(f"slide {slide.get('slide_id', '?')} rendered_image", resolved, target_ratio)
                source_images.append(item)
                if issue:
                    issues.append(issue)

    if render_manifest:
        manifest = _load_json(render_manifest)
        entries = manifest.get("slides", [])
        if len(entries) != report["slide_count"]:
            issues.append(
                f"render_manifest count {len(entries)} does not match pptx slide_count {report['slide_count']}"
            )
        for entry in entries:
            if entry.get("status") != "completed":
                issues.append(f"slide {entry.get('slide_id', '?')}: render status is {entry.get('status')}")
            for key in ("prompt_file", "backend", "generated_source", "copied_to"):
                if not entry.get(key):
                    issues.append(f"slide {entry.get('slide_id', '?')}: missing manifest field {key}")

    if deck_json:
        deck = _load_json(deck_json)
        slides = deck.get("slides", [])
        if len(slides) != report["slide_count"]:
            issues.append(f"deck_json count {len(slides)} does not match pptx slide_count {report['slide_count']}")
        for idx, slide in enumerate(slides, start=1):
            if not slide.get("background"):
                issues.append(f"deck_json slide {idx}: missing background")
                continue
            resolved = _resolve(deck_json.parent, slide.get("background"))
            if resolved and resolved not in checked_paths:
                checked_paths.add(resolved)
                item, issue = _aspect_issue(f"deck_json slide {idx} background", resolved, target_ratio)
                source_images.append(item)
                if issue:
                    issues.append(issue)

    report["status"] = "fail" if issues else "pass"
    report["issues"] = issues
    report["source_images"] = source_images
    return report


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--spec", type=Path)
    parser.add_argument("--render-manifest", type=Path)
    parser.add_argument("--deck-json", type=Path)
    parser.add_argument("--min-slides", type=int, default=1)
    parser.add_argument("--out", type=Path)
    args = parser.parse_args()

    report = audit_package(
        args.pptx,
        args.spec,
        args.render_manifest,
        args.deck_json,
        args.min_slides,
    )
    text = json.dumps(report, ensure_ascii=False, indent=2)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(text, encoding="utf-8")
    print(text)
    if report["status"] == "fail":
        raise SystemExit(1)


if __name__ == "__main__":
    main()
